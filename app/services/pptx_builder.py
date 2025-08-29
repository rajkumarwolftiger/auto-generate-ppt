from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def build_presentation(template_path, slides, notes, out_path):
    tpl = Presentation(template_path)
    prs = Presentation(template_path)  # start from template to inherit masters

    # Simple approach: use first "Title and Content" layout if available
    title_content_layout = None
    for i, layout in enumerate(prs.slide_layouts):
        # crude check by name; you can inspect placeholders for robustness
        if "title" in layout.name.lower() and "content" in layout.name.lower():
            title_content_layout = layout
            break
    if title_content_layout is None:
        title_content_layout = prs.slide_layouts[0]

    for idx, s in enumerate(slides):
        slide = prs.slides.add_slide(title_content_layout)
        # title
        slide.shapes.title.text = s["title"][:120] if slide.shapes.title else s["title"]

        # body
        body = slide.placeholders[1].text_frame if len(slide.placeholders) > 1 else None
        if body:
            body.clear()
            for j, bullet in enumerate(s["bullets"][:8]):
                p = body.add_paragraph() if j else body.paragraphs[0]
                p.text = bullet
                p.level = 0

        # reuse an image from the template’s first slide if present
        # (You can get fancier by cycling images or matching layouts)
        if len(tpl.slides) > 0:
            for shp in tpl.slides[0].shapes:
                if shp.shape_type == 13:  # PICTURE
                    # add at same position/size
                    slide.shapes.add_picture(shp.image.blob, shp.left, shp.top, width=shp.width, height=shp.height)
                    break

        # optional speaker notes
        if notes and idx < len(notes):
            slide.notes_slide.notes_text_frame.text = notes[idx]

    prs.save(out_path)
